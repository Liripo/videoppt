import numpy as np
import cv2 as cv
import os
from PIL import ImageGrab
import sys
from skimage.measure import compare_ssim
#
import argparse
import pptx
from pptx.util import Inches

parser = argparse.ArgumentParser(description='该程序可进行对屏幕进行录屏，\
    分析视频帧数相似性。')
parser.add_argument('-r',action = "store_true",help='使用录屏程序,按q退出录制,结果保存为video.avi。')
parser.add_argument('-c',action = "store_true",
                    help='是否使用相机录制，默认使用截图录制,\
                    你可以添加这个参数启用,请确保有内置摄像头。')
parser.add_argument('-f',
                    help='输入分析视频文件名。')
parser.add_argument('--ahash',action = "store_true",
                    help = '使用ahash算法。')

parser.add_argument('--same',action = "store_true",
                    help = '使用same,完全一样,这可能会把所有的帧数取出。\
                    如果视频中ppt完全不怎么变化可以使用。')                    
parser.add_argument('-s',type = float,default = 0.98,
                    help="去除高于此相似度的图片。")
                    
parser.add_argument('--frameskip',type = int,default = 0,
                    help='跳帧')
parser.add_argument('-p',action = "store_true",help = "使用生成ppt程序。")
parser.add_argument('-d',
                    help='输入需要整理的文件夹路径，\
                    会导入文件夹下所有jpg的图片。')
parser.add_argument('--ssim',action = "store_true",
                    help = '使用ssim算法。')                
args = parser.parse_args()
#--------------------function

def cap_image(carmer_bool):
  #录屏
  #
  #---------------------
  #获取屏幕大小
  print(carmer_bool)
  print("屏幕录制中，按q结束！！")
  curScreen = ImageGrab.grab()
  height, width = curScreen.size
  #帧数
  fps = 20
  video = cv.VideoWriter('video.avi', 
    cv.VideoWriter_fourcc(*'XVID'), fps, (height, width))
  if(carmer_bool == True):
    cap = cv.VideoCapture(0)
    while(True):
      ret, frame = cap.read()
      gray = cv.cvtColor(frame, cv.COLOR_RGB2BGR)
      cv.imshow('frame',gray)
      video.write(gray)
      k = cv.waitKey(1000)
      if k == ord('q'):
        break
    cap.release()
    cv.destroyAllWindows()
  if carmer_bool == False:
    imageNum = 0
    while True:
      imageNum += 1
      captureImage = ImageGrab.grab()  
      frame = cv.cvtColor(np.array(captureImage), cv.COLOR_RGB2BGR)
      # 显示无图像的窗口
      cv.imshow('capturing', np.zeros((1, 255), np.uint8))
      # 控制窗口显示位置，方便通过按键方式退出
      cv.moveWindow('capturing', height - 100, width - 100)
      #延时录制参数可加，if imageNum > fps * start:
      video.write(frame)
      if cv.waitKey(1000) & 0xFF == ord('q'):
        break
#---api算法
def compare_same(frame1,frame2):
  difference = cv.subtract(frame1,frame2)
  result = not np.any(difference)
  if result:
    similarity = 1
  else:
    similarity = 0
  return similarity
  
#----ahash算法
def ahash(img):
  img = cv.resize(img,(8,8))
  gray = cv.cvtColor(img, cv.COLOR_BGR2GRAY)
  s = 0
  hash_str = ""
  #总像素，暂时先这样实现吧，应该有内置总像素值
  for i in range(8):
    for j in range(8):
      s = s+gray[i, j]
  #平均值
  avg = s/64
  for i in range(8):
    for j in range(8):
      if gray[i,j] > avg :
        hash_str = hash_str + "1"
      else:
        hash_str = hash_str + "0"
  return hash_str

##--------------------------------------比较函数
def ahash_similarity(frame1,frame2):
  h1 = ahash(frame1)
  h2 = ahash(frame2)
  #相似度使用汉明距离
  n = 0
  if len(h1) != len(h2):
    print("哈希值不同，请排除bug!")
    sys.exit(1)
  for i in range(len(h1)):
    if h1[i] != h2[i]:
      n = n + 1
  similarity = 1 - (n/len(h1))
  return similarity
def compare_ssim2(frame1,frame2):
  grayA = cv.cvtColor(frame1, cv.COLOR_BGR2GRAY)
  grayB = cv.cvtColor(frame2, cv.COLOR_BGR2GRAY)

  # compute the Structural Similarity Index (SSIM) between the two
  # images, ensuring that the difference image is returned
  (score, diff) = compare_ssim(grayA, grayB, full=True)
  diff = (diff * 255).astype("uint8")
  #-----------print("SSIM: {}".format(score))
  return score
#-----------------------------------
def run(file,std_similarity,compare_fun,frameskip):
  cap = cv.VideoCapture(file)
  frame_count = int(cap.get(cv.CAP_PROP_FRAME_COUNT))
  print("总计%d帧" % frame_count)
  idx = 0
  while (idx < (frame_count - 1)):
    cap.set(1, idx)
    ret1,frame1 = cap.read()
    cv.imwrite(f'./result/frame{idx}.jpg',frame1)
    #cv.imshow("image",frame) #show
    if (idx + frameskip) > frame_count:
      cv.imwrite(f'./result/frame{idx}.jpg',frame1)
      break
    for idy in range(idx + frameskip,frame_count):
      print(idx,"-",idy,"<<总帧数=",frame_count,">>")
      cap.set(1,idy)
      ret2,frame2 = cap.read()
      #相似度
      similarity = compare_fun(frame1,frame2)
      print("相似度",similarity)
      if similarity < std_similarity and idy != (frame_count - 1):
        print("分析到第%d帧数" % idy)
        print("此处帧与上个帧相似度为%d"%similarity,"开始跳至下个帧。")
        idx = idy
        break
        #最后退出的情况
      if idy == (frame_count - 1):
        idx = idy
        break
        
        
def produce_ppt(dir_jpg):
  pptFile = pptx.Presentation()
  picFiles = [fn for fn in os.listdir(dir_jpg) if fn.endswith('.jpg')]
  #文件按数字排序
  picFiles.sort(key = lambda x: int(x[5:-4]))
  for jpg in picFiles:
    #加入空白页布局
    slide = pptFile.slides.add_slide(pptFile.slide_layouts[6])
    jpg = dir_jpg+ "/" + jpg
    print(jpg)
    width = Inches(8)
    height = Inches(6.7)
    left = Inches(0.5)
    top = Inches(0.5)
    slide.shapes.add_picture(jpg, 
      left,top,width,height)
  pptFile.save('zzz.pptx')
    
##############
#main##########
################

if __name__ == "__main__":
  print("工作目录=>",os.getcwd())
  if(args.r):
    cap_image(args.c)
  if(args.f is not None):
    if not os.path.exists("result"):
      os.makedirs("result")
    print("分析文件=>%s" % args.f)
    if (args.ahash == True):
      run(args.f,std_similarity = args.s,compare_fun = ahash_similarity,
        frameskip = args.frameskip)
    if (args.same == True):
      run(args.f,std_similarity = args.s,compare_fun = compare_same,
      frameskip = args.frameskip)
    if (args.ssim == True):
      run(args.f,std_similarity = args.s,compare_fun = compare_ssim2,
      frameskip = args.frameskip)
  if(args.p is not None):
    produce_ppt(args.d)
      
  
  
